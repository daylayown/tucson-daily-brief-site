#!/usr/bin/env python3
"""Build the UA JOUR 428 guest lecture presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Theme colors ---
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)       # deep navy
BG_ACCENT = RGBColor(0x22, 0x22, 0x3A)      # slightly lighter
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
TEXT_DIM = RGBColor(0x99, 0x99, 0xAA)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x4F)  # warm amber
ACCENT_BLUE = RGBColor(0x6C, 0xA0, 0xDC)    # soft blue

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=TEXT_WHITE, bold=False,
             space_before=Pt(6), alignment=PP_ALIGN.LEFT, font_name="Arial"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_divider(slide, top, color=ACCENT_ORANGE, width_inches=2):
    left = Inches(1)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(width_inches), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
         "Building Things in Journalism", font_size=48, bold=True,
         color=ACCENT_ORANGE)

add_text(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.6),
         "Nicholas De Leon", font_size=28, color=TEXT_WHITE)

tf = add_text(slide, Inches(1), Inches(4.1), Inches(11), Inches(1),
              "Senior Reporter, Consumer Reports", font_size=18, color=TEXT_DIM)
add_para(tf, "JOUR 428: Product Development in Journalism", font_size=18,
         color=TEXT_DIM)
add_para(tf, "University of Arizona  \u2022  March 5, 2026", font_size=18,
         color=TEXT_DIM)

# =============================================================================
# SLIDE 2 — Bio / Career Path
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "20 years in journalism", font_size=36, bold=True,
         color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

tf = add_text(slide, Inches(1), Inches(2.4), Inches(10.5), Inches(3.5),
              "Started at Gizmodo in December 2005.",
              font_size=24, color=TEXT_WHITE)
add_para(tf, "", font_size=12, color=TEXT_DIM)
add_para(tf, "Gizmodo \u2192 TechCrunch \u2192 The Daily (News Corp) \u2192 Circa \u2192 Vice \u2192 Consumer Reports",
         font_size=20, color=TEXT_DIM)
add_para(tf, "", font_size=12, color=TEXT_DIM)
add_para(tf, "At Consumer Reports since 2017, covering laptops, PCs, WiFi, and AI.",
         font_size=24, color=TEXT_WHITE)
add_para(tf, "", font_size=12, color=TEXT_DIM)
add_para(tf, "NYU \u2014 Journalism & Political Science, 2008",
         font_size=18, color=TEXT_DIM)
add_para(tf, "Moved from NYC to Catalina, AZ in 2023.",
         font_size=18, color=TEXT_DIM)

# =============================================================================
# SLIDE 3 — Unorthodox Inspirations
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Where I actually learned journalism instincts",
         font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

inspirations = [
    ("Tech blogs (2003\u20132008)",
     "Engadget, Gizmodo, Ars Technica \u2014 fast, opinionated, obsessive about getting it right. Proved you didn\u2019t need a masthead to matter."),
    ("Wrestling Observer Newsletter",
     "Dave Meltzer has published a niche newsletter since 1983. Decades before Substack, he proved one expert voice + a dedicated audience = a sustainable business."),
    ("Leo Laporte / TWiT",
     "Built a podcast network before \u201cpodcast\u201d was a word. Showed me that audio + niche expertise + consistency = audience."),
]

y = Inches(2.2)
for title, desc in inspirations:
    add_text(slide, Inches(1), y, Inches(10.5), Inches(0.4),
             title, font_size=22, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1), y + Inches(0.45), Inches(10.5), Inches(0.8),
             desc, font_size=16, color=TEXT_DIM)
    y += Inches(1.4)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "The point: journalism instincts can come from anywhere. Follow your obsessions.",
         font_size=18, color=TEXT_WHITE, bold=True)

# =============================================================================
# SLIDE 4 — What is Tucson Daily Brief?
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Tucson Daily Brief", font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

tf = add_text(slide, Inches(1), Inches(2.2), Inches(5.5), Inches(4),
              "An AI-powered local news briefing for Tucson, updated daily at 6 AM.",
              font_size=22, color=TEXT_WHITE)
add_para(tf, "", font_size=12, color=TEXT_DIM)
add_para(tf, "\u2022  Website \u2014 tucsondailybrief.com", font_size=18, color=TEXT_DIM)
add_para(tf, "\u2022  Podcast \u2014 Apple Podcasts, Spotify", font_size=18, color=TEXT_DIM)
add_para(tf, "\u2022  YouTube \u2014 daily video briefings", font_size=18, color=TEXT_DIM)
add_para(tf, "\u2022  Telegram \u2014 morning delivery", font_size=18, color=TEXT_DIM)
add_para(tf, "", font_size=12, color=TEXT_DIM)
add_para(tf, "No human touches it. One person built it.", font_size=20,
         color=TEXT_WHITE, bold=True)

# Placeholder for screenshot
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2), Inches(4.8), Inches(4.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = BG_ACCENT
shape.line.color.rgb = TEXT_DIM
shape.line.width = Pt(1)
tf_ph = shape.text_frame
tf_ph.word_wrap = True
p = tf_ph.paragraphs[0]
p.text = "[ screenshot of tucsondailybrief.com ]"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_DIM
p.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 5 — How I Built It (Architecture)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "How it works", font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

steps = [
    ("1", "Curate sources", "20+ local RSS feeds, organized by editorial priority"),
    ("2", "Write the playbook", "AI instructions: what to cover, how to rank, what tone to use"),
    ("3", "AI generates briefing", "Claude reads the feeds, writes the briefing following editorial rules"),
    ("4", "Distribute everywhere", "Website \u2192 Podcast \u2192 YouTube \u2192 Telegram \u2014 one pipeline"),
]

y = Inches(2.2)
for num, title, desc in steps:
    # Number circle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), y, Inches(0.6), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    shape.line.fill.background()
    tf_n = shape.text_frame
    tf_n.paragraphs[0].text = num
    tf_n.paragraphs[0].font.size = Pt(20)
    tf_n.paragraphs[0].font.color.rgb = BG_DARK
    tf_n.paragraphs[0].font.bold = True
    tf_n.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, Inches(2), y, Inches(9), Inches(0.4),
             title, font_size=22, color=TEXT_WHITE, bold=True)
    add_text(slide, Inches(2), y + Inches(0.4), Inches(9), Inches(0.4),
             desc, font_size=16, color=TEXT_DIM)
    y += Inches(1.1)

add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
         "Total: ~300 lines of Python  \u2022  Zero frameworks  \u2022  Free hosting on GitHub Pages",
         font_size=16, color=TEXT_DIM)

# =============================================================================
# SLIDE 6 — Live Demo
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
         "Live demo", font_size=54, bold=True,
         color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "Let\u2019s look at the code.",
         font_size=28, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 7 — Newsletters That Work
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Newsletters that work", font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

newsletters = [
    ("Wrestling Observer Newsletter",
     "Dave Meltzer, since 1983",
     "One expert, one niche, 40+ years. The original Substack before Substack existed."),
    ("Platformer",
     "Casey Newton",
     "Left The Verge to cover Big Tech independently. Broke major stories as a solo operation."),
    ("404 Media",
     "Jason Koebler, Emanuel Maiberg, Samantha Cole, Joseph Cox",
     "Four journalists left Vice and built their own newsroom. Profitable within months."),
]

y = Inches(2.2)
for name, byline, desc in newsletters:
    add_text(slide, Inches(1), y, Inches(10.5), Inches(0.4),
             name, font_size=22, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1), y + Inches(0.4), Inches(10.5), Inches(0.35),
             byline, font_size=15, color=TEXT_DIM)
    add_text(slide, Inches(1), y + Inches(0.8), Inches(10.5), Inches(0.6),
             desc, font_size=17, color=TEXT_WHITE)
    y += Inches(1.45)

# =============================================================================
# SLIDE 8 — What Makes Newsletters Succeed
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "What makes them work", font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

principles = [
    ("Niche focus",
     "Don\u2019t try to be everything. Be the best at one thing. The narrower your focus, the more indispensable you are."),
    ("Consistent voice",
     "People subscribe to people, not brands. Your perspective is your product."),
    ("Sustainable cadence",
     "Daily, weekly, whatever \u2014 but never miss. Reliability builds trust faster than quality alone."),
    ("Direct relationship",
     "Email, RSS, podcast \u2014 you own the audience. No algorithm deciding who sees your work."),
]

y = Inches(2.2)
for title, desc in principles:
    add_text(slide, Inches(1), y, Inches(10.5), Inches(0.4),
             title, font_size=22, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1), y + Inches(0.4), Inches(10.5), Inches(0.7),
             desc, font_size=17, color=TEXT_DIM)
    y += Inches(1.15)

# =============================================================================
# SLIDE 9 — Careers in 2026
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Careers in 2026", font_size=36, bold=True, color=ACCENT_ORANGE)
add_divider(slide, Inches(1.6))

skills = [
    ("Think clearly",
     "Every tool is useless without judgment. Know what matters and why."),
    ("Write well",
     "AI can generate text. It can\u2019t decide what\u2019s worth saying."),
    ("Build things",
     "You don\u2019t need to be an engineer. But knowing how to make something real \u2014 a site, a feed, a product \u2014 changes what\u2019s possible."),
    ("Stay curious",
     "The tools will keep changing. The instinct to learn won\u2019t become obsolete."),
]

y = Inches(2.2)
for title, desc in skills:
    add_text(slide, Inches(1), y, Inches(10.5), Inches(0.4),
             title, font_size=22, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1), y + Inches(0.4), Inches(10.5), Inches(0.7),
             desc, font_size=17, color=TEXT_DIM)
    y += Inches(1.1)

add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "Skills matter more than job titles.",
         font_size=20, color=TEXT_WHITE, bold=True)

# =============================================================================
# SLIDE 10 — The Next Gizmodo Is In Your Head
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
         "The next Gizmodo is in your head.",
         font_size=48, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(2.5),
         "You don\u2019t need permission to start.\n"
         "You don\u2019t need a newsroom to publish.\n"
         "You don\u2019t need venture capital to reach people.",
         font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.7),
         "You need a topic you care about, a voice, and the willingness to ship.",
         font_size=22, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 11 — Q&A
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
         "Questions?", font_size=54, bold=True,
         color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 12 — Thank You / Contact
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "Thanks!", font_size=48, bold=True, color=ACCENT_ORANGE,
         alignment=PP_ALIGN.CENTER)

contacts = [
    "nicholas.deleon@consumer.org",
    "linkedin.com/in/nicholas-de-leon-3b5b6a9",
    "tucsondailybrief.com",
]

y = Inches(3.3)
for c in contacts:
    add_text(slide, Inches(1), y, Inches(11), Inches(0.5),
             c, font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    y += Inches(0.55)

add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Nicholas De Leon  \u2022  Senior Reporter, Consumer Reports",
         font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# --- Save ---
out_path = "/home/nicholas/claude-code-projects/tucson-daily-brief-site/ua_jour428_presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
